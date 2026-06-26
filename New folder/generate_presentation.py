from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_bullets_slide(prs, title, bullets, level_map=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()

    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = b
        p.level = 0 if not level_map else level_map.get(i, 0)
        p.font.size = Pt(20)

    return slide


def add_project_summary(prs):
    bullets = [
        "Background: MSICT and DFHQ relied on manual paper logbooks for visitor registration.",
        "Problem Statement: Manual logs caused long queues, difficult record retrieval, low data security, and frequent errors.",
        "Objectives (Measurable):",
        "1) Reduce average registration time by digitizing check-in and check-out.",
        "2) Maintain 100% searchable digital visitor records.",
        "3) Generate daily, weekly, and monthly reports automatically.",
        "Scope: Covers user login, visitor registration, entry/exit tracking, visitor history, and reporting at MSICT/DFHQ.",
    ]
    level_map = {3: 1, 4: 1, 5: 1}
    return add_bullets_slide(prs, "1. Project Summary", bullets, level_map)


def add_literature_review(prs):
    bullets = [
        "Information Systems Theory supports using ICT to improve organizational efficiency (Laudon & Laudon, 2019).",
        "Access Control Theory emphasizes regulating entry to protect resources (Sandhu et al., 1996).",
        "Empirical findings show digital VMS improves identification, speed, and traceability (ATT Systems Group, 2025; IJMRSET, 2025).",
        "National digital transformation efforts encourage computerized systems in Tanzania (Tanzania ICT Authority, 2023).",
        "Research gap: Existing manual systems at MSICT/DFHQ lacked secure centralized tracking and automated reporting.",
    ]
    return add_bullets_slide(prs, "2. Literature Review", bullets)


def add_methodology(prs):
    bullets = [
        "Research Design: Descriptive + exploratory with mixed methods.",
        "Data Collection Techniques:",
        "- Questionnaires (visitors and staff)",
        "- Interviews (security officers, administrators, ICT personnel)",
        "- Observation (registration workflow at entry points)",
        "- Document review (logbooks and institutional records)",
        "System Development Methodology: Waterfall Model.",
        "Waterfall Stages: Planning -> Requirement Analysis -> Design -> Implementation -> Testing -> Deployment -> Maintenance.",
    ]
    level_map = {2: 1, 3: 1, 4: 1, 5: 1}
    return add_bullets_slide(prs, "3. Methodology", bullets, level_map)


def add_timeline_gantt(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "4. Project Timeline (Gantt Chart)"

    left = Inches(0.5)
    top = Inches(1.4)
    width = Inches(12.3)
    height = Inches(5.3)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(245, 247, 250)
    bg.line.color.rgb = RGBColor(210, 215, 222)

    tasks = [
        ("Requirement Analysis", 1, 2),
        ("System Design", 3, 5),
        ("Coding", 6, 10),
        ("Testing", 11, 11),
        ("Deployment", 12, 12),
    ]

    start_x = Inches(4.2)
    end_x = Inches(12.2)
    y0 = Inches(2.0)
    row_h = Inches(0.7)
    bar_h = Inches(0.35)

    total_weeks = 12
    step = (end_x - start_x) / total_weeks

    # Week headers and vertical grid
    for w in range(1, total_weeks + 1):
        x = start_x + step * (w - 1)
        label = slide.shapes.add_textbox(x, Inches(1.55), step, Inches(0.3))
        label_tf = label.text_frame
        label_tf.text = f"W{w}"
        label_tf.paragraphs[0].font.size = Pt(11)
        label_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y0 - Inches(0.15), Pt(1), Inches(4.3))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(225, 230, 236)
        line.line.fill.background()

    colors = [
        RGBColor(53, 132, 228),
        RGBColor(63, 185, 80),
        RGBColor(246, 153, 63),
        RGBColor(214, 90, 90),
        RGBColor(137, 99, 186),
    ]

    for i, (task, ws, we) in enumerate(tasks):
        y = y0 + row_h * i
        # Task label
        tbox = slide.shapes.add_textbox(Inches(0.8), y, Inches(3.3), Inches(0.4))
        tf = tbox.text_frame
        tf.text = task
        tf.paragraphs[0].font.size = Pt(14)

        bx = start_x + step * (ws - 1)
        bw = step * (we - ws + 1)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, y + Inches(0.05), bw, bar_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors[i % len(colors)]
        bar.line.color.rgb = RGBColor(255, 255, 255)

    note = slide.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.4))
    ntf = note.text_frame
    ntf.text = "Timeline reflects current project milestones from the report (12-week plan)."
    ntf.paragraphs[0].font.size = Pt(12)
    ntf.paragraphs[0].font.italic = True

    return slide


def add_references(prs):
    bullets = [
        "ATT Systems Group. (2025). Automated visitor management solutions for modern organizations.",
        "IJMRSET. (2025). A study on automated visitor management systems, 7(2), 45-52.",
        "Laudon, K. C., & Laudon, J. P. (2019). Management information systems (15th ed.). Pearson.",
        "Sandhu, R. S., Coyne, E. J., Feinstein, H. L., & Youman, C. E. (1996). Role-based access control models. IEEE Computer, 29(2), 38-47.",
        "Tanzania ICT Authority. (2023). ICT systems adoption and digital transformation in Tanzania.",
    ]
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "References"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = b
        p.level = 0
        p.font.size = Pt(14)

    return slide


def style_titles(prs):
    for slide in prs.slides:
        if slide.shapes.title is not None:
            t = slide.shapes.title.text_frame.paragraphs[0]
            t.font.size = Pt(34)
            t.font.bold = True
            t.font.color.rgb = RGBColor(16, 55, 92)


def add_footer(prs):
    for slide in prs.slides:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.25))
        tf = box.text_frame
        tf.text = "Prepared in line with Progress Report Presentation Requirements - The Research Office"
        p = tf.paragraphs[0]
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(90, 98, 108)
        p.alignment = PP_ALIGN.RIGHT


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Visitor Management System",
        "Progress Report Presentation\nMSICT and DFHQ",
    )
    add_project_summary(prs)
    add_literature_review(prs)
    add_methodology(prs)
    add_timeline_gantt(prs)
    add_references(prs)
    style_titles(prs)
    add_footer(prs)

    out1 = r"c:\\Users\\hp\\Desktop\\RIPOTI_modified_presentation.pptx"
    out2 = r"c:\\xampp\\htdocs\\Visitors_MS\\New folder\\RIPOTI_modified_presentation.pptx"
    prs.save(out1)
    prs.save(out2)
    print(out1)
    print(out2)


if __name__ == "__main__":
    main()
